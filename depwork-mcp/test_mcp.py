#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""WPS MCP Server 全面端到端测试。

模拟 AI 通过 MCP 协议完整操作 WPS Office 的流程：
1. 创建文档 → 2. 添加内容 → 3. 导出文件 → 4. 验证内容
"""

import sys
import os
import asyncio
import json
import subprocess

sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from wps_controller.mcp_server import server

WORKDIR = os.path.dirname(os.path.abspath(__file__))
passed = 0
failed = 0


def check(label, ok, detail=""):
    global passed, failed
    icon = "✅" if ok else "❌"
    if ok:
        passed += 1
    else:
        failed += 1
    print(f"  {icon} {label}" + (f" → {detail}" if detail else ""))


def file_exists(path):
    return os.path.exists(path) and os.path.getsize(path) > 0


async def call_tool(name: str, args: dict = None) -> dict:
    """调用 MCP 工具，返回解析后的 JSON 数据。"""
    try:
        result = await server.call_tool(name, args or {})
    except Exception as e:
        return {"success": False, "error": str(e)}
    if result.is_error:
        return {"success": False, "error": str(result.content)}
    if result.content and hasattr(result.content[0], 'text'):
        return json.loads(result.content[0].text)
    return {"success": False, "error": "no content"}


def verify_docx(path, expect_paras, expect_table_data=None, expect_images=0):
    from docx import Document
    doc = Document(path)
    errors = []
    actual_paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for exp in expect_paras:
        if not any(exp in p for p in actual_paras):
            errors.append(f"段落缺失: '{exp}'")
    if expect_table_data:
        if not doc.tables:
            errors.append("表格不存在")
        else:
            t = doc.tables[0]
            actual_rows = [[c.text.strip() for c in row.cells] for row in t.rows]
            for exp_row in expect_table_data:
                found = any(all(any(e in cell for cell in r) for e in exp_row) for r in actual_rows)
                if not found:
                    errors.append(f"表格行缺失: {exp_row}")
    if expect_images > 0:
        if len(doc.inline_shapes) < expect_images:
            errors.append(f"图片: 期望{expect_images}, 实际{len(doc.inline_shapes)}")
    return errors


def verify_xlsx(path, expect_cells, expect_sheets=None, sheet_name=None):
    from openpyxl import load_workbook
    wb = load_workbook(path)
    errors = []
    if expect_sheets:
        for s in expect_sheets:
            if s not in wb.sheetnames:
                errors.append(f"工作表缺失: '{s}'")
    ws = wb[sheet_name] if sheet_name else wb.worksheets[0]
    if expect_cells:
        for ref, val in expect_cells.items():
            actual = ws[ref].value
            if actual is None or str(val) not in str(actual):
                errors.append(f"{ref}: 期望'{val}', 实际'{actual}'")
    return errors


def verify_pptx(path, expect_slide_count=None, expect_titles=None):
    from pptx import Presentation
    prs = Presentation(path)
    errors = []
    if expect_slide_count is not None:
        if len(prs.slides) != expect_slide_count:
            errors.append(f"幻灯片数: 期望{expect_slide_count}, 实际{len(prs.slides)}")
    if expect_titles:
        for i, title in enumerate(expect_titles):
            if i >= len(prs.slides):
                errors.append(f"幻灯片{i}不存在")
                continue
            slide = prs.slides[i]
            found = False
            for shape in slide.placeholders:
                if shape.has_text_frame and title in shape.text_frame.text:
                    found = True
                    break
            if not found:
                all_text = " ".join(
                    shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
                )
                if title not in all_text:
                    errors.append(f"幻灯片{i}: 期望'{title}', 实际'{all_text[:50]}'")
    return errors


async def main():
    print("=" * 70)
    print("  🚀 WPS MCP Server 全面端到端测试")
    print("=" * 70)

    # 清理 WPS 进程
    subprocess.run(["taskkill", "/F", "/IM", "wps.exe", "/T"], capture_output=True)
    subprocess.run(["taskkill", "/F", "/IM", "et.exe", "/T"], capture_output=True)
    subprocess.run(["taskkill", "/F", "/IM", "wpp.exe", "/T"], capture_output=True)

    # ══════════════════════════════════════════════════════════════
    # 1. Writer 全流程
    # ══════════════════════════════════════════════════════════════
    print("\n  📝 Writer 全流程（MCP → 创建 → 编辑 → 导出 → 验证）")

    r = await call_tool("create_document", {"doc_type": "writer", "name": "MCP Writer 测试"})
    check("创建 Writer 文档", r["success"], r.get("message", ""))

    r = await call_tool("writer_add_heading", {"text": "MCP 驱动的 WPS 自动化", "level": 1})
    check("添加 H1 标题", r["success"])

    r = await call_tool("writer_add_paragraph", {
        "text": "本文档由 AI 通过 MCP 协议自动生成，展示端到端控制能力。"
    })
    check("添加段落", r["success"])

    r = await call_tool("writer_add_list", {
        "items": ["MCP 协议传输", "内存会话管理", "COM 自动化渲染"],
        "list_style": "bullet"
    })
    check("添加列表", r["success"])

    r = await call_tool("writer_add_table", {
        "rows": 4, "cols": 3,
        "data": [
            ["模块", "工具数", "状态"],
            ["Writer", 9, "✅"],
            ["Calc", 6, "✅"],
            ["Impress", 6, "✅"],
        ]
    })
    check("添加表格 4x3", r["success"])

    r = await call_tool("writer_add_image", {
        "image_path": os.path.join(WORKDIR, "test_image.png"),
        "width": "8cm", "height": "6cm"
    })
    check("添加图片", r["success"])

    r = await call_tool("get_document_info", {})
    check("查看文档信息", r["success"] and r["data"]["content_count"] == 5,
          f"content_count={r['data']['content_count']}")

    # 撤销测试
    r = await call_tool("undo", {})
    check("撤销图片", r["success"])

    r = await call_tool("get_document_info", {})
    check("撤销后内容数=4", r["success"] and r["data"]["content_count"] == 4)

    # 重做
    r = await call_tool("redo", {})
    check("重做图片", r["success"])

    r = await call_tool("get_document_info", {})
    check("重做后内容数=5", r["success"] and r["data"]["content_count"] == 5)

    # 查找替换
    r = await call_tool("writer_find_replace", {"find_text": "MCP", "replace_text": "MCP(Model Context Protocol)"})
    check("查找替换", r["success"], f"replaced={r['data']['replaced']}")

    # 快速导出
    subprocess.run(["taskkill", "/F", "/IM", "wps.exe", "/T"], capture_output=True)
    docx_path = os.path.join(WORKDIR, "mcp_writer_test.docx")
    r = await call_tool("export_document", {
        "output_path": docx_path, "preset": "docx", "overwrite": True
    })
    check("快速导出 DOCX", r["success"],
          f"{r['data']['file_size']:,} bytes" if r["success"] else r.get("error", ""))

    if file_exists(docx_path):
        errors = verify_docx(docx_path,
            expect_paras=["MCP", "AI", "COM"],
            expect_table_data=[["Writer", "✅"], ["Calc", "✅"], ["Impress", "✅"]],
            expect_images=1)
        check("验证 DOCX 内容", len(errors) == 0, "; ".join(errors) if errors else "段落+表格+图片 正确")

    # Live 导出
    subprocess.run(["taskkill", "/F", "/IM", "wps.exe", "/T"], capture_output=True)
    docx_live = os.path.join(WORKDIR, "mcp_writer_live.docx")
    r = await call_tool("export_document", {
        "output_path": docx_live, "preset": "docx", "overwrite": True,
        "live": True, "speed": "fast"
    })
    check("Live 导出 DOCX (fast)", r["success"],
          f"{r['data']['file_size']:,} bytes, method={r['data'].get('method','')}" if r["success"] else r.get("error",""))

    if file_exists(docx_live):
        errors = verify_docx(docx_live,
            expect_paras=["MCP", "AI"],
            expect_table_data=[["Writer", "✅"]],
            expect_images=1)
        check("验证 Live DOCX 内容", len(errors) == 0, "; ".join(errors) if errors else "正确")

    # ══════════════════════════════════════════════════════════════
    # 2. Calc 全流程
    # ══════════════════════════════════════════════════════════════
    print("\n  📊 Calc 全流程（MCP → 创建 → 编辑 → 导出 → 验证）")

    r = await call_tool("create_document", {"doc_type": "calc", "name": "MCP Calc 测试"})
    check("创建 Calc 文档", r["success"])

    r = await call_tool("calc_set_cell", {"ref": "A1", "value": "产品"})
    check("设置 A1", r["success"])
    r = await call_tool("calc_set_cell", {"ref": "B1", "value": "销量"})
    check("设置 B1", r["success"])
    r = await call_tool("calc_set_cell", {"ref": "C1", "value": "单价"})
    check("设置 C1", r["success"])

    r = await call_tool("calc_set_range", {
        "start_ref": "A2",
        "data": [["苹果", 100, 5], ["香蕉", 200, 3], ["橙子", 150, 4]]
    })
    check("批量写入 A2:C4", r["success"], f"cells={r['data']['cells_set']}")

    r = await call_tool("calc_set_cell", {"ref": "D1", "value": "总额"})
    r = await call_tool("calc_set_cell", {
        "ref": "D2", "value": "", "formula": "=B2*C2"
    })
    check("设置公式 D2=B2*C2", r["success"])

    r = await call_tool("calc_merge_cells", {"start_ref": "A1", "end_ref": "D1"})
    check("合并 A1:D1", r["success"])

    r = await call_tool("calc_add_sheet", {"name": "汇总"})
    check("添加工作表[汇总]", r["success"])

    r = await call_tool("calc_list_sheets", {})
    check("列出工作表", r["success"] and len(r["data"]) == 2)

    # 快速导出
    subprocess.run(["taskkill", "/F", "/IM", "et.exe", "/T"], capture_output=True)
    xlsx_path = os.path.join(WORKDIR, "mcp_calc_test.xlsx")
    r = await call_tool("export_document", {
        "output_path": xlsx_path, "preset": "xlsx", "overwrite": True
    })
    check("快速导出 XLSX", r["success"],
          f"{r['data']['file_size']:,} bytes" if r["success"] else r.get("error",""))

    if file_exists(xlsx_path):
        errors = verify_xlsx(xlsx_path,
            expect_cells={"A1": "产品", "A2": "苹果", "B2": "100"},
            expect_sheets=["Sheet1", "汇总"],
            sheet_name="Sheet1")
        check("验证 XLSX 内容", len(errors) == 0, "; ".join(errors) if errors else "正确")

    # Live 导出
    subprocess.run(["taskkill", "/F", "/IM", "et.exe", "/T"], capture_output=True)
    xlsx_live = os.path.join(WORKDIR, "mcp_calc_live.xlsx")
    r = await call_tool("export_document", {
        "output_path": xlsx_live, "preset": "xlsx", "overwrite": True,
        "live": True, "speed": "fast"
    })
    check("Live 导出 XLSX (fast)", r["success"],
          f"{r['data']['file_size']:,} bytes" if r["success"] else r.get("error",""))

    if file_exists(xlsx_live):
        errors = verify_xlsx(xlsx_live,
            expect_cells={"A1": "产品", "A2": "苹果"},
            expect_sheets=["Sheet1", "汇总"],
            sheet_name="Sheet1")
        check("验证 Live XLSX 内容", len(errors) == 0, "; ".join(errors) if errors else "正确")

    # ══════════════════════════════════════════════════════════════
    # 3. Impress 全流程
    # ══════════════════════════════════════════════════════════════
    print("\n  🎬 Impress 全流程（MCP → 创建 → 编辑 → 导出 → 验证）")

    r = await call_tool("create_document", {"doc_type": "impress", "name": "MCP Impress 测试"})
    check("创建 Impress 文档", r["success"])

    r = await call_tool("impress_add_slide", {"title": "MCP 协议", "content": "模型上下文协议"})
    check("添加幻灯片1", r["success"])
    r = await call_tool("impress_add_slide", {"title": "架构", "content": "AI → MCP → WPS"})
    check("添加幻灯片2", r["success"])
    r = await call_tool("impress_add_slide", {"title": "演示", "content": "实时打字效果"})
    check("添加幻灯片3", r["success"])

    r = await call_tool("impress_add_element", {
        "slide_index": 0, "element_type": "text_box",
        "text": "Powered by MCP", "x": "3cm", "y": "5cm",
        "width": "8cm", "height": "3cm"
    })
    check("添加文本框元素", r["success"])

    r = await call_tool("impress_duplicate_slide", {"index": 1})
    check("复制幻灯片1", r["success"])

    r = await call_tool("impress_list_slides", {})
    check("列出幻灯片", r["success"] and len(r["data"]) == 4, f"slides={len(r['data'])}")

    # 快速导出
    subprocess.run(["taskkill", "/F", "/IM", "wpp.exe", "/T"], capture_output=True)
    pptx_path = os.path.join(WORKDIR, "mcp_impress_test.pptx")
    r = await call_tool("export_document", {
        "output_path": pptx_path, "preset": "pptx", "overwrite": True
    })
    check("快速导出 PPTX", r["success"],
          f"{r['data']['file_size']:,} bytes" if r["success"] else r.get("error",""))

    if file_exists(pptx_path):
        errors = verify_pptx(pptx_path, expect_slide_count=4,
            expect_titles=["MCP 协议", "架构", "架构"])
        check("验证 PPTX 内容", len(errors) == 0, "; ".join(errors) if errors else "正确")

    # Live 导出
    subprocess.run(["taskkill", "/F", "/IM", "wpp.exe", "/T"], capture_output=True)
    pptx_live = os.path.join(WORKDIR, "mcp_impress_live.pptx")
    r = await call_tool("export_document", {
        "output_path": pptx_live, "preset": "pptx", "overwrite": True,
        "live": True, "speed": "fast"
    })
    check("Live 导出 PPTX (fast)", r["success"],
          f"{r['data']['file_size']:,} bytes" if r["success"] else r.get("error",""))

    if file_exists(pptx_live):
        errors = verify_pptx(pptx_live, expect_slide_count=4,
            expect_titles=["MCP 协议", "架构"])
        check("验证 Live PPTX 内容", len(errors) == 0, "; ".join(errors) if errors else "正确")

    # ══════════════════════════════════════════════════════════════
    # 4. 辅助工具测试
    # ══════════════════════════════════════════════════════════════
    print("\n  🔧 辅助工具测试")

    r = await call_tool("list_export_presets", {})
    check("列出导出预设", r["success"] and len(r["data"]) >= 10, f"{len(r['data'])} 个预设")

    r = await call_tool("list_page_profiles", {})
    check("列出页面配置", r["success"] and len(r["data"]) >= 3, f"{len(r['data'])} 个配置")

    r = await call_tool("get_session_status", {})
    check("会话状态", r["success"] and r["data"]["has_project"] is True)

    # 保存到文件
    save_path = os.path.join(WORKDIR, "mcp_saved_project.json")
    r = await call_tool("save_document", {"path": save_path})
    check("保存项目到文件", r["success"], r.get("message", ""))

    r = await call_tool("open_document", {"path": save_path})
    check("重新打开项目", r["success"], r.get("message", ""))

    # ══════════════════════════════════════════════════════════════
    # 5. 汇总
    # ══════════════════════════════════════════════════════════════
    print("\n" + "=" * 70)
    print("  📋 MCP Server 全面测试汇总")
    print("=" * 70)
    print(f"  总用例: {passed + failed}")
    print(f"  通过:   {passed}")
    print(f"  失败:   {failed}")
    print(f"  通过率: {passed / (passed + failed) * 100:.1f}%")

    print("\n  生成的文件:")
    for f in ["mcp_writer_test.docx", "mcp_writer_live.docx",
              "mcp_calc_test.xlsx", "mcp_calc_live.xlsx",
              "mcp_impress_test.pptx", "mcp_impress_live.pptx"]:
        fp = os.path.join(WORKDIR, f)
        if file_exists(fp):
            print(f"    ✅ {f} ({os.path.getsize(fp):,} bytes)")
        else:
            print(f"    ❌ {f}")

    print("\n  MCP 工具统计:")
    tools = await server.list_tools()
    print(f"    总工具数: {len(tools)}")
    cats = {"writer": 0, "calc": 0, "impress": 0, "export": 0, "document": 0, "session": 0, "list": 0}
    for t in tools:
        name = t.name.lower()
        if "writer" in name:
            cats["writer"] += 1
        elif "calc" in name:
            cats["calc"] += 1
        elif "impress" in name:
            cats["impress"] += 1
        elif "export" in name:
            cats["export"] += 1
        elif "create" in name or "open" in name or "save" in name or "info" in name:
            cats["document"] += 1
        elif "undo" in name or "redo" in name or "status" in name:
            cats["session"] += 1
        elif "list" in name:
            cats["list"] += 1
    for cat, count in cats.items():
        print(f"    {cat}: {count} 个工具")

    print("\n" + "=" * 70)
    if failed == 0:
        print("  🎉 全部通过！MCP Server 可投入生产使用。")
    else:
        print("  ⚠️  有失败项，请检查。")
    print("=" * 70)


if __name__ == "__main__":
    asyncio.run(main())
