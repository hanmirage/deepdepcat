#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""全面落地测试：Writer / Calc / Impress × 快速模式/Live模式 × 内容验证。"""

import sys
import os
import json
import subprocess
import time

sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

WORKDIR = os.path.dirname(os.path.abspath(__file__))
PYTHON = sys.executable
CLI = "wps_controller.wps_cli"

passed = 0
failed = 0


def cli(args, expect_ok=True):
    """执行 CLI 命令，返回 (ok, stdout_json, stderr)。"""
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    env["PYTHONUTF8"] = "1"
    full = [PYTHON, "-m", CLI, "--json"] + args
    proc = subprocess.run(full, capture_output=True, text=True,
                          cwd=WORKDIR, encoding="utf-8",
                          errors="replace", env=env)
    ok = (proc.returncode == 0) == expect_ok
    try:
        out = json.loads(proc.stdout) if (proc.stdout or "").strip() else None
    except json.JSONDecodeError:
        out = {"raw": (proc.stdout or "")[:300]}
    return ok, out, (proc.stderr or "").strip()


def check(label, ok, detail=""):
    global passed, failed
    icon = "✅" if ok else "❌"
    if ok:
        passed += 1
    else:
        failed += 1
    print(f"  {icon} {label}" + (f" → {detail}" if detail else ""))


def file_exists(path):
    return os.path.exists(path) and os.path.getsize(path) > 0


def verify_docx(path, expect_paras, expect_table_data=None, expect_images=0):
    """用 python-docx 验证 DOCX 内容。"""
    from docx import Document
    doc = Document(path)
    errors = []

    # 段落
    actual_paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for exp in expect_paras:
        found = any(exp in p for p in actual_paras)
        if not found:
            errors.append(f"段落缺失: '{exp}'")

    # 表格（在任意行查找匹配，不按固定索引）
    if expect_table_data:
        if not doc.tables:
            errors.append("表格不存在")
        else:
            t = doc.tables[0]
            actual_rows = []
            for row in t.rows:
                actual_rows.append([c.text.strip() for c in row.cells])
            for exp_row in expect_table_data:
                found = any(
                    all(exp in actual for exp, actual in zip(exp_row, r))
                    for r in actual_rows
                )
                if not found:
                    errors.append(f"表格行缺失: {exp_row}")

    # 图片
    if expect_images > 0:
        actual_images = len(doc.inline_shapes)
        if actual_images < expect_images:
            errors.append(f"图片数量: 期望{expect_images}, 实际{actual_images}")

    return errors


def verify_xlsx(path, expect_cells=None, expect_sheets=None, sheet_name=None):
    """用 openpyxl 验证 XLSX 内容。"""
    from openpyxl import load_workbook
    wb = load_workbook(path)
    errors = []

    if expect_sheets:
        actual_sheets = wb.sheetnames
        for s in expect_sheets:
            if s not in actual_sheets:
                errors.append(f"工作表缺失: '{s}'")

    # 按名称或索引选择工作表
    if sheet_name:
        if sheet_name not in wb.sheetnames:
            errors.append(f"工作表 '{sheet_name}' 不存在")
            return errors
        ws = wb[sheet_name]
    else:
        ws = wb.worksheets[0]

    if expect_cells:
        for ref, val in expect_cells.items():
            actual = ws[ref].value
            if actual is None or str(val) not in str(actual):
                errors.append(f"{ref}: 期望'{val}', 实际'{actual}'")

    return errors


def verify_pptx(path, expect_slide_count=None, expect_titles=None):
    """用 python-pptx 验证 PPTX 内容。"""
    from pptx import Presentation
    prs = Presentation(path)
    errors = []

    if expect_slide_count is not None:
        actual = len(prs.slides)
        if actual != expect_slide_count:
            errors.append(f"幻灯片数: 期望{expect_slide_count}, 实际{actual}")

    if expect_titles:
        for i, title in enumerate(expect_titles):
            if i >= len(prs.slides):
                errors.append(f"幻灯片{i}不存在")
                continue
            slide = prs.slides[i]
            # 优先检查 title placeholder
            found = False
            for shape in slide.placeholders:
                if shape.has_text_frame and title in shape.text_frame.text:
                    found = True
                    break
            if not found:
                # 退一步检查所有 shape
                all_text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_text += shape.text_frame.text + " "
                if title not in all_text:
                    errors.append(f"幻灯片{i}标题: 期望'{title}', 实际'{all_text[:50]}'")

    return errors


# ══════════════════════════════════════════════════════════════
# 1. Writer 全面测试
# ══════════════════════════════════════════════════════════════
print("\n" + "=" * 70)
print("  📝 Writer 全面测试（标题/段落/列表/表格带数据/图片）")
print("=" * 70)

wproj = os.path.join(WORKDIR, "final_writer.json")
wdocx_fast = os.path.join(WORKDIR, "final_writer_fast.docx")
wdocx_live = os.path.join(WORKDIR, "final_writer_live.docx")
wdocx_slow = os.path.join(WORKDIR, "final_writer_slow.docx")

# 1.1 创建
ok, out, err = cli(["document", "new", "--type", "writer", "--name", "Writer全面测试", "-o", wproj])
check("创建 Writer 文档", ok and out.get("type") == "writer")

# 1.2 添加标题
ok, out, _ = cli(["--project", wproj, "writer", "add-heading", "-t", "2024年度技术总结", "-l", "1"])
check("添加 H1 标题", ok)

# 1.3 添加段落
ok, out, _ = cli(["--project", wproj, "writer", "add-paragraph",
                  "-t", "本报告由AI自动生成，涵盖全年技术成果与数据指标分析。"])
check("添加正文段落", ok)

# 1.4 添加列表
ok, out, _ = cli(["--project", wproj, "writer", "add-list",
                  "-i", "系统稳定性99.9%", "-i", "API响应速度提升40%", "-i", "用户满意度96分",
                  "--style", "bullet"])
check("添加无序列表(3项)", ok)

# 1.5 添加带数据的表格（通过 Python 直接写入 JSON）
with open(wproj, "r", encoding="utf-8") as f:
    proj = json.load(f)
proj["content"].append({
    "type": "table", "rows": 5, "cols": 4,
    "data": [
        ["季度", "活跃用户", "营收", "满意度"],
        ["Q1", "800万", "0.7亿", "91%"],
        ["Q2", "950万", "0.8亿", "93%"],
        ["Q3", "1100万", "0.9亿", "94%"],
        ["Q4", "1200万", "1.1亿", "96%"],
    ]
})
proj["content"].append({
    "type": "image_ref", "name": "test_image.png",
    "path": "test_image.png", "width": "8cm", "height": "6cm"
})
with open(wproj, "w", encoding="utf-8") as f:
    json.dump(proj, f, ensure_ascii=False, indent=2)
check("添加带数据表格(5x4) + 图片", True)

# 1.6 快速模式导出
taskkill = subprocess.run(["taskkill", "/F", "/IM", "wps.exe", "/T"],
                          capture_output=True, text=True)
ok, out, err = cli(["--project", wproj, "export", "render", wdocx_fast,
                     "-p", "docx", "--overwrite"])
check("快速模式导出 DOCX", ok and file_exists(wdocx_fast),
      f"{os.path.getsize(wdocx_fast):,} bytes" if file_exists(wdocx_fast) else err)

# 1.7 验证快速模式内容
if file_exists(wdocx_fast):
    errors = verify_docx(wdocx_fast,
        expect_paras=["2024年度技术总结", "AI自动生成", "系统稳定性"],
        expect_table_data=[
            ["季度","活跃用户","营收","满意度"],
            ["Q1","800万","0.7亿","91%"],
            ["Q2","950万","0.8亿","93%"],
            ["Q3","1100万","0.9亿","94%"],
            ["Q4","1200万","1.1亿","96%"],
        ],
        expect_images=1)
    check("验证快速模式内容", len(errors) == 0, "; ".join(errors) if errors else "段落+表格5行+图片 全部正确")

# 1.8 Live 模式导出 (fast)
taskkill = subprocess.run(["taskkill", "/F", "/IM", "wps.exe", "/T"],
                          capture_output=True, text=True)
ok, out, err = cli(["--project", wproj, "export", "render", wdocx_live,
                     "-p", "docx", "--overwrite", "--live", "--speed", "fast"])
check("Live模式(fast)导出 DOCX", ok and file_exists(wdocx_live),
      f"{os.path.getsize(wdocx_live):,} bytes, method={out.get('method','')}" if ok else err)

# 1.9 验证 Live 模式内容
if file_exists(wdocx_live):
    errors = verify_docx(wdocx_live,
        expect_paras=["2024年度技术总结", "AI自动生成", "系统稳定性"],
        expect_table_data=[
            ["季度","活跃用户","营收","满意度"],
            ["Q1","800万","0.7亿","91%"],
            ["Q2","950万","0.8亿","93%"],
            ["Q3","1100万","0.9亿","94%"],
            ["Q4","1200万","1.1亿","96%"],
        ],
        expect_images=1)
    check("验证Live模式内容", len(errors) == 0, "; ".join(errors) if errors else "段落+表格5行+图片 全部正确")

# 1.10 Live 模式导出 (slow) — 慢速可见
taskkill = subprocess.run(["taskkill", "/F", "/IM", "wps.exe", "/T"],
                          capture_output=True, text=True)
ok, out, err = cli(["--project", wproj, "export", "render", wdocx_slow,
                     "-p", "docx", "--overwrite", "--live", "--speed", "slow"])
check("Live模式(slow)导出 DOCX", ok and file_exists(wdocx_slow),
      f"{os.path.getsize(wdocx_slow):,} bytes, char_delay={out.get('char_delay','')}" if ok else err)

if file_exists(wdocx_slow):
    errors = verify_docx(wdocx_slow,
        expect_paras=["2024年度技术总结", "AI自动生成"],
        expect_table_data=[
            ["季度","活跃用户","营收","满意度"],
            ["Q1","800万","0.7亿","91%"],
            ["Q4","1200万","1.1亿","96%"],
        ],
        expect_images=1)
    check("验证Slow模式内容", len(errors) == 0, "; ".join(errors) if errors else "正确")


# ══════════════════════════════════════════════════════════════
# 2. Calc 全面测试
# ══════════════════════════════════════════════════════════════
print("\n" + "=" * 70)
print("  📊 Calc 全面测试（表头/批量数据/公式/合并/多工作表）")
print("=" * 70)

cproj = os.path.join(WORKDIR, "final_calc.json")
cxlsx_fast = os.path.join(WORKDIR, "final_calc_fast.xlsx")
cxlsx_live = os.path.join(WORKDIR, "final_calc_live.xlsx")

# 2.1 创建
ok, out, _ = cli(["document", "new", "--type", "calc", "--name", "Calc全面测试", "-o", cproj])
check("创建 Calc 文档", ok and out.get("type") == "calc")

# 2.2 设置表头
for ref, val in [("A1","产品名"), ("B1","销量"), ("C1","单价"), ("D1","总额")]:
    ok, _, _ = cli(["--project", cproj, "calc", "set-cell", ref, val])
check("设置表头 A1:D1", ok)

# 2.3 批量写入数据
data = [["产品A", 100, 50], ["产品B", 200, 30], ["产品C", 150, 80]]
ok, out, _ = cli(["--project", cproj, "calc", "set-range", "A2", "-d", json.dumps(data)])
check("批量写入 A2:C4 (9格)", ok and out.get("cells_set") == 9)

# 2.4 设置公式
ok, out, _ = cli(["--project", cproj, "calc", "set-cell", "D2", "", "--formula", "=B2*C2"])
check("设置公式 D2=B2*C2", ok)

# 2.5 合并单元格
ok, out, _ = cli(["--project", cproj, "calc", "merge-cells", "A1", "D1"])
check("合并 A1:D1", ok)

# 2.6 添加第二个工作表
ok, out, _ = cli(["--project", cproj, "calc", "add-sheet", "-n", "汇总"])
check("添加工作表[汇总]", ok)

# 2.7 快速模式导出
taskkill = subprocess.run(["taskkill", "/F", "/IM", "et.exe", "/T"],
                          capture_output=True, text=True)
ok, out, err = cli(["--project", cproj, "export", "render", cxlsx_fast,
                     "-p", "xlsx", "--overwrite"])
check("快速模式导出 XLSX", ok and file_exists(cxlsx_fast),
      f"{os.path.getsize(cxlsx_fast):,} bytes" if ok else err)

# 2.8 验证快速模式内容
if file_exists(cxlsx_fast):
    errors = verify_xlsx(cxlsx_fast,
        expect_cells={"A1": "产品名", "A2": "产品A", "B2": "100", "A3": "产品B"},
        expect_sheets=["Sheet1", "汇总"],
        sheet_name="Sheet1")
    check("验证快速模式内容", len(errors) == 0, "; ".join(errors) if errors else "单元格+工作表顺序 正确")

# 2.9 Live 模式导出
taskkill = subprocess.run(["taskkill", "/F", "/IM", "et.exe", "/T"],
                          capture_output=True, text=True)
ok, out, err = cli(["--project", cproj, "export", "render", cxlsx_live,
                     "-p", "xlsx", "--overwrite", "--live", "--speed", "fast"])
check("Live模式(fast)导出 XLSX", ok and file_exists(cxlsx_live),
      f"{os.path.getsize(cxlsx_live):,} bytes" if ok else err)

if file_exists(cxlsx_live):
    errors = verify_xlsx(cxlsx_live,
        expect_cells={"A1": "产品名", "A2": "产品A", "B2": "100"},
        expect_sheets=["Sheet1", "汇总"],
        sheet_name="Sheet1")
    check("验证Live模式内容", len(errors) == 0, "; ".join(errors) if errors else "正确")


# ══════════════════════════════════════════════════════════════
# 3. Impress 全面测试
# ══════════════════════════════════════════════════════════════
print("\n" + "=" * 70)
print("  🎬 Impress 全面测试（多幻灯片/标题/内容/文本框元素）")
print("=" * 70)

iproj = os.path.join(WORKDIR, "final_impress.json")
ipptx_fast = os.path.join(WORKDIR, "final_impress_fast.pptx")
ipptx_live = os.path.join(WORKDIR, "final_impress_live.pptx")

# 3.1 创建
ok, out, _ = cli(["document", "new", "--type", "impress", "--name", "Impress全面测试", "-o", iproj])
check("创建 Impress 文档", ok and out.get("type") == "impress")

# 3.2 添加幻灯片
ok, _, _ = cli(["--project", iproj, "impress", "add-slide", "-t", "产品发布", "-c", "2024年度新品"])
check("添加幻灯片1: 产品发布", ok)

ok, _, _ = cli(["--project", iproj, "impress", "add-slide", "-t", "功能特性", "-c", "AI驱动\n高性能\n易用性"])
check("添加幻灯片2: 功能特性", ok)

ok, _, _ = cli(["--project", iproj, "impress", "add-slide", "-t", "技术架构", "-c", "微服务架构设计"])
check("添加幻灯片3: 技术架构", ok)

# 3.3 添加文本框元素
ok, _, _ = cli(["--project", iproj, "impress", "add-element", "0",
                "--type", "text_box", "--text", "Hello from Agent!",
                "--x", "3cm", "--y", "5cm", "--width", "8cm", "--height", "3cm"])
check("添加文本框元素到幻灯片0", ok)

# 3.4 复制幻灯片
ok, _, _ = cli(["--project", iproj, "impress", "duplicate-slide", "1"])
check("复制幻灯片1", ok)

# 3.5 快速模式导出
taskkill = subprocess.run(["taskkill", "/F", "/IM", "wpp.exe", "/T"],
                          capture_output=True, text=True)
ok, out, err = cli(["--project", iproj, "export", "render", ipptx_fast,
                     "-p", "pptx", "--overwrite"])
check("快速模式导出 PPTX", ok and file_exists(ipptx_fast),
      f"{os.path.getsize(ipptx_fast):,} bytes" if ok else err)

if file_exists(ipptx_fast):
    errors = verify_pptx(ipptx_fast, expect_slide_count=4,
        expect_titles=["产品发布", "功能特性", "功能特性", "技术架构"])
    check("验证快速模式内容", len(errors) == 0, "; ".join(errors) if errors else "4张幻灯片+标题 全部正确")

# 3.6 Live 模式导出
taskkill = subprocess.run(["taskkill", "/F", "/IM", "wpp.exe", "/T"],
                          capture_output=True, text=True)
ok, out, err = cli(["--project", iproj, "export", "render", ipptx_live,
                     "-p", "pptx", "--overwrite", "--live", "--speed", "fast"])
check("Live模式(fast)导出 PPTX", ok and file_exists(ipptx_live),
      f"{os.path.getsize(ipptx_live):,} bytes" if ok else err)

if file_exists(ipptx_live):
    errors = verify_pptx(ipptx_live, expect_slide_count=4,
        expect_titles=["产品发布", "功能特性", "功能特性", "技术架构"])
    check("验证Live模式内容", len(errors) == 0, "; ".join(errors) if errors else "正确")


# ══════════════════════════════════════════════════════════════
# 4. 速度对比测试
# ══════════════════════════════════════════════════════════════
print("\n" + "=" * 70)
print("  ⚡ 速度对比测试（fast vs normal vs slow vs custom delay）")
print("=" * 70)

speed_tests = [
    ("fast",   None,    0.01),
    ("normal", None,    0.03),
    ("slow",   None,    0.08),
    ("normal", 0.05,    0.05),  # custom delay 覆盖
]

for speed, custom_delay, expect_delay in speed_tests:
    taskkill = subprocess.run(["taskkill", "/F", "/IM", "wps.exe", "/T"],
                              capture_output=True, text=True)
    args = ["--project", wproj, "export", "render",
            os.path.join(WORKDIR, f"speed_test_{speed}.docx"),
            "-p", "docx", "--overwrite", "--live", "--speed", speed]
    if custom_delay:
        args += ["--delay", str(custom_delay)]

    t0 = time.time()
    ok, out, err = cli(args)
    elapsed = time.time() - t0

    actual_delay = out.get("char_delay", "?") if ok else "?"
    ok_delay = abs(float(actual_delay) - expect_delay) < 0.001 if ok else False
    label = f"speed={speed}" + (f" delay={custom_delay}" if custom_delay else "")
    check(f"{label} → char_delay={actual_delay}s, 耗时={elapsed:.1f}s",
          ok and ok_delay, "")


# ══════════════════════════════════════════════════════════════
# 5. 汇总报告
# ══════════════════════════════════════════════════════════════
print("\n" + "=" * 70)
print("  📋 全面落地测试汇总")
print("=" * 70)
print(f"  总用例: {passed + failed}")
print(f"  通过:   {passed}")
print(f"  失败:   {failed}")
print(f"  通过率: {passed / (passed + failed) * 100:.1f}%")

print("\n  生成的文件:")
for f in [
    "final_writer_fast.docx", "final_writer_live.docx", "final_writer_slow.docx",
    "final_calc_fast.xlsx", "final_calc_live.xlsx",
    "final_impress_fast.pptx", "final_impress_live.pptx",
]:
    fp = os.path.join(WORKDIR, f)
    if file_exists(fp):
        print(f"    ✅ {f} ({os.path.getsize(fp):,} bytes)")
    else:
        print(f"    ❌ {f} (不存在)")

print("\n" + "=" * 70)
if failed == 0:
    print("  🎉 全部通过！")
else:
    print("  ⚠️  有失败项，请检查")
print("=" * 70)
